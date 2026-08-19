from pathlib import Path

from openai import AsyncOpenAI

from docx import Document
from docx.shared import Pt, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT

from pptx import Presentation
from pptx.util import Inches, Pt as PPTPt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .config import OPENAI_API_KEY, OPENAI_MODEL


Path("output").mkdir(exist_ok=True)

client = (
    AsyncOpenAI(api_key=OPENAI_API_KEY)
    if OPENAI_API_KEY
    else None
)


# =========================================================
# AI
# =========================================================

async def generate(
    kind,
    topic,
    pages=None,
    template=None,
    design=None,
):
    if not client:
        return "OPENAI_API_KEY sozlanmagan."

    extra = ""

    if pages:
        extra += f"""
Hujjat hajmi: {pages} bet.
Hajmni mazmunni takrorlash yoki keraksiz cho'zish orqali emas,
mazmunli va to'liq ilmiy material orqali ta'minla.
"""

    if template:
        extra += f"""
Tanlangan shablon: {template}.
"""

    if design:
        extra += f"""
Tanlangan slayd dizayni: {design}.
"""

    kind_lower = kind.lower()

    if kind_lower == "kurs ishi":
        instruction = """
KURS ISHI TAYYORLA.

Tuzilishi:

1. TITUL VARAQ
2. MUNDARIJA
3. KIRISH
4. I BOB
   1.1.
   1.2.
   1.3.
5. II BOB
   2.1.
   2.2.
   2.3.
6. XULOSA
7. FOYDALANILGAN ADABIYOTLAR

Kirishda:
- mavzuning dolzarbligi;
- tadqiqot maqsadi;
- tadqiqot vazifalari;
- tadqiqot obyekti;
- tadqiqot predmeti
yoritilsin.

Har bir bobda 3 ta mazmunli bo'lim bo'lsin.

Ilmiy uslubdan foydalan.
Mavzuni to'liq ochib ber.
Bir xil gaplarni takrorlama.
"""

    elif kind_lower == "mustaqil ish":
        instruction = """
MUSTAQIL ISH TAYYORLA.

Tuzilishi:

1. TITUL VARAQ
2. REJA
3. KIRISH
4. ASOSIY QISM
5. XULOSA
6. FOYDALANILGAN ADABIYOTLAR

Mavzu to'liq va izchil yoritilsin.
"""

    elif kind_lower == "referat":
        instruction = """
REFERAT TAYYORLA.

Tuzilishi:

1. TITUL VARAQ
2. REJA
3. KIRISH
4. ASOSIY QISM
5. XULOSA
6. FOYDALANILGAN ADABIYOTLAR

Ilmiy va tushunarli uslubdan foydalan.
"""

    elif kind_lower == "maqola":
        instruction = """
ILMIY MAQOLA TAYYORLA.

Tuzilishi:

Sarlavha
Muallif
Annotatsiya
Kalit so'zlar
Kirish
Adabiyotlar tahlili
Tadqiqot metodologiyasi
Natijalar va muhokama
Xulosa
Foydalanilgan adabiyotlar

Maqola ilmiy maqola talablariga mos,
mantiqiy va tabiiy yozilsin.
"""

    elif kind_lower == "tezis":
        instruction = """
ILMIY TEZIS TAYYORLA.

Tuzilishi:

Sarlavha
Muallif
Asosiy mazmun
Ilmiy natijalar
Xulosa
Kalit so'zlar
Foydalanilgan manbalar
"""

    elif kind_lower == "konspekt":
        instruction = """
MAVZU BO'YICHA KONSPEKT TAYYORLA.

Muhim tushunchalar,
asosiy fikrlar,
faktlar,
ta'riflar
va xulosalarni tartibli joylashtir.
"""

    elif kind_lower == "glossariy":
        instruction = """
MAVZU BO'YICHA GLOSSARIY TAYYORLA.

Muhim ilmiy atamalarni tanla.
Har bir atamaga aniq, qisqa va ilmiy ta'rif ber.
"""

    else:
        instruction = """
Vazifani professional, tabiiy va ravon
o'zbek tilida bajar.
"""

    prompt = f"""
Sen professional o'zbek tilidagi
ta'lim va ilmiy ishlar AI yordamchisisan.

Vazifa turi:
{kind}

Mavzu:
{topic}

{extra}

{instruction}

Umumiy talablar:

- O'zbek tilida yoz.
- Grammatik jihatdan to'g'ri yoz.
- Tabiiy inson yozganidek bo'lsin.
- Bir xil fikrlarni takrorlama.
- Keraksiz AI izohlarini yozma.
- Sarlavhalarni aniq ajrat.
- Mavzudan chetga chiqma.
- Imkon qadar mazmunli va batafsil yoz.
"""

    response = await client.chat.completions.create(
        model=OPENAI_MODEL,
        messages=[
            {
                "role": "system",
                "content": (
                    "Sen professional o'zbek tilidagi "
                    "ta'lim va ilmiy ishlar AI yordamchisisan."
                ),
            },
            {
                "role": "user",
                "content": prompt,
            },
        ],
        temperature=0.6,
    )

    return response.choices[0].message.content or ""


async def answer(q):
    if not client:
        return "OPENAI_API_KEY sozlanmagan."

    response = await client.chat.completions.create(
        model=OPENAI_MODEL,
        messages=[
            {
                "role": "user",
                "content": q,
            }
        ],
        temperature=0.5,
    )

    return response.choices[0].message.content or ""


# =========================================================
# WORD
# =========================================================

def setup_document(doc):

    section = doc.sections[0]

    # Hoshiyalar
    section.top_margin = Cm(2)
    section.bottom_margin = Cm(2)
    section.left_margin = Cm(3)
    section.right_margin = Cm(1.5)

    style = doc.styles["Normal"]

    style.font.name = "Times New Roman"
    style.font.size = Pt(14)

    # Times New Roman
    rpr = style._element.get_or_add_rPr()
    rfonts = rpr.rFonts

    if rfonts is not None:
        rfonts.set(
            "{http://schemas.openxmlformats.org/"
            "wordprocessingml/2006/main}ascii",
            "Times New Roman",
        )
        rfonts.set(
            "{http://schemas.openxmlformats.org/"
            "wordprocessingml/2006/main}hAnsi",
            "Times New Roman",
        )

    paragraph = style.paragraph_format

    paragraph.line_spacing = 1.5
    paragraph.first_line_indent = Cm(1.25)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)


def add_text_paragraph(doc, text):

    if not text.strip():
        return

    p = doc.add_paragraph()

    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.first_line_indent = Cm(1.25)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.space_after = Pt(0)

    run = p.add_run(text.strip())

    run.font.name = "Times New Roman"
    run.font.size = Pt(14)

    return p


def add_heading(doc, text):

    p = doc.add_paragraph()

    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    p.paragraph_format.line_spacing = 1.5
    p.paragraph_format.first_line_indent = Cm(0)
    p.paragraph_format.space_before = Pt(8)
    p.paragraph_format.space_after = Pt(8)

    run = p.add_run(text.strip())

    run.bold = True
    run.font.name = "Times New Roman"
    run.font.size = Pt(14)

    return p


# =========================================================
# TITUL
# =========================================================

def add_title_page(doc, title, kind):

    p = doc.add_paragraph()

    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    r = p.add_run(
        "GULISTON DAVLAT UNIVERSITETI"
    )

    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(14)

    for _ in range(2):
        doc.add_paragraph()

    p = doc.add_paragraph()

    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    r = p.add_run(kind.upper())

    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(16)

    for _ in range(2):
        doc.add_paragraph()

    p = doc.add_paragraph()

    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    r = p.add_run(title)

    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(16)

    for _ in range(5):
        doc.add_paragraph()

    p = doc.add_paragraph()

    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    r = p.add_run(
        "Bajardi: ______________________________\n"
        "Tekshirdi: _____________________________"
    )

    r.font.name = "Times New Roman"
    r.font.size = Pt(14)

    for _ in range(4):
        doc.add_paragraph()

    p = doc.add_paragraph()

    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    r = p.add_run("Guliston — 2026")

    r.font.name = "Times New Roman"
    r.font.size = Pt(14)

    doc.add_page_break()


# =========================================================
# DOCX
# =========================================================

def make_docx(
    title,
    body,
    kind="hujjat",
    pages=None,
):

    safe_title = (
        title[:60]
        .replace("/", "_")
        .replace("\\", "_")
        .replace(":", "_")
        .replace("*", "_")
        .replace("?", "_")
        .replace('"', "_")
        .replace("<", "_")
        .replace(">", "_")
        .replace("|", "_")
    )

    path = Path("output") / f"{safe_title}.docx"

    doc = Document()

    setup_document(doc)

    # Titul
    if kind.lower() in [
        "kurs ishi",
        "mustaqil ish",
        "referat",
    ]:
        add_title_page(
            doc,
            title,
            kind,
        )

    lines = body.splitlines()

    for line in lines:

        text = line.strip()

        if not text:
            continue

        upper = text.upper()

        is_heading = (
            upper in [
                "KIRISH",
                "XULOSA",
                "MUNDARIJA",
                "ASOSIY QISM",
                "FOYDALANILGAN ADABIYOTLAR",
            ]
            or upper.startswith("I BOB")
            or upper.startswith("II BOB")
            or upper.startswith("1.")
            or upper.startswith("2.")
        )

        if is_heading:
            add_heading(
                doc,
                text,
            )
        else:
            add_text_paragraph(
                doc,
                text,
            )

    doc.save(path)

    return path


# =========================================================
# SLAYD DIZAYNLARI
# =========================================================

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


DESIGNS = {

    "🎓 Akademik": {
        "font": "Times New Roman",
        "title_size": 30,
        "body_size": 21,
        "bg": RGBColor(245, 248, 252),
        "primary": RGBColor(25, 72, 125),
        "secondary": RGBColor(67, 114, 170),
        "text": RGBColor(35, 45, 55),
        "white": RGBColor(255, 255, 255),
    },

    "💼 Professional": {
        "font": "Arial",
        "title_size": 30,
        "body_size": 21,
        "bg": RGBColor(248, 248, 248),
        "primary": RGBColor(35, 42, 50),
        "secondary": RGBColor(90, 100, 110),
        "text": RGBColor(45, 45, 45),
        "white": RGBColor(255, 255, 255),
    },

    "✨ Zamonaviy": {
        "font": "Aptos",
        "title_size": 32,
        "body_size": 22,
        "bg": RGBColor(28, 32, 48),
        "primary": RGBColor(230, 170, 70),
        "secondary": RGBColor(80, 110, 180),
        "text": RGBColor(245, 245, 248),
        "white": RGBColor(255, 255, 255),
    },

    "🧊 Minimal": {
        "font": "Arial",
        "title_size": 29,
        "body_size": 20,
        "bg": RGBColor(255, 255, 255),
        "primary": RGBColor(40, 40, 40),
        "secondary": RGBColor(180, 180, 180),
        "text": RGBColor(65, 65, 65),
        "white": RGBColor(255, 255, 255),
    },

    "🌈 Kreativ": {
        "font": "Aptos",
        "title_size": 32,
        "body_size": 22,
        "bg": RGBColor(248, 244, 255),
        "primary": RGBColor(105, 65, 170),
        "secondary": RGBColor(220, 90, 160),
        "text": RGBColor(55, 45, 70),
        "white": RGBColor(255, 255, 255),
    },
}


# =========================================================
# YORDAMCHI FUNKSIYALAR
# =========================================================

def set_slide_background(slide, color):
    """
    Slayd fonini belgilaydi.
    """

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(
    slide,
    shape_type,
    left,
    top,
    width,
    height,
    color,
):
    """
    Slaydga dekorativ shakl qo'shadi.
    """

    shape = slide.shapes.add_shape(
        shape_type,
        left,
        top,
        width,
        height,
    )

    shape.fill.solid()
    shape.fill.fore_color.rgb = color

    shape.line.fill.background()

    return shape


def add_title(
    slide,
    text,
    font,
    size,
    color,
    dark=False,
):
    """
    Slayd sarlavhasini yaratadi.
    """

    box = slide.shapes.add_textbox(
        Inches(0.75),
        Inches(0.45),
        Inches(11.8),
        Inches(1.0),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    paragraph = frame.paragraphs[0]

    paragraph.alignment = PP_ALIGN.CENTER

    run = paragraph.add_run()
    run.text = text

    run.font.name = font
    run.font.size = PPTPt(size)
    run.font.bold = True
    run.font.color.rgb = color

    return box


def add_body(
    slide,
    lines,
    font,
    size,
    color,
):
    """
    Slayd asosiy matnini yaratadi.
    """

    box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(1.75),
        Inches(11.25),
        Inches(4.75),
    )

    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True

    for index, text in enumerate(lines):

        paragraph = (
            frame.paragraphs[0]
            if index == 0
            else frame.add_paragraph()
        )

        paragraph.text = text
        paragraph.level = 0
        paragraph.space_after = PPTPt(12)

        for run in paragraph.runs:
            run.font.name = font
            run.font.size = PPTPt(size)
            run.font.color.rgb = color

    return box


def add_slide_number(
    slide,
    number,
    font,
    color,
):
    """
    Slayd raqamini chiqaradi.
    """

    box = slide.shapes.add_textbox(
        Inches(11.9),
        Inches(6.75),
        Inches(0.7),
        Inches(0.35),
    )

    frame = box.text_frame
    frame.clear()

    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT

    run = paragraph.add_run()
    run.text = str(number)

    run.font.name = font
    run.font.size = PPTPt(12)
    run.font.color.rgb = color


# =========================================================
# AKADEMIK
# =========================================================

def design_academic(
    slide,
    title,
    lines,
    number,
    style,
):
    """
    Rasmiy universitet/ilmiy uslub.
    """

    set_slide_background(
        slide,
        style["bg"],
    )

    # Yuqori ko'k chiziq
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.18),
        style["primary"],
    )

    # Chap dekorativ blok
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0.18),
        Inches(0.18),
        Inches(7.32),
        style["secondary"],
    )

    add_title(
        slide,
        title,
        style["font"],
        style["title_size"],
        style["primary"],
    )

    add_body(
        slide,
        lines,
        style["font"],
        style["body_size"],
        style["text"],
    )

    # Pastki chiziq
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.75),
        Inches(6.55),
        Inches(11.6),
        Inches(0.035),
        style["secondary"],
    )

    add_slide_number(
        slide,
        number,
        style["font"],
        style["primary"],
    )


# =========================================================
# PROFESSIONAL
# =========================================================

def design_professional(
    slide,
    title,
    lines,
    number,
    style,
):
    """
    Jiddiy va biznes uslubidagi dizayn.
    """

    set_slide_background(
        slide,
        style["bg"],
    )

    # Chap qora panel
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.28),
        Inches(7.5),
        style["primary"],
    )

    # Sarlavha tagidagi chiziq
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        Inches(1.38),
        Inches(11.6),
        Inches(0.06),
        style["primary"],
    )

    add_title(
        slide,
        title,
        style["font"],
        style["title_size"],
        style["primary"],
    )

    add_body(
        slide,
        lines,
        style["font"],
        style["body_size"],
        style["text"],
    )

    # Pastki kichik dekor
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        Inches(6.65),
        Inches(1.4),
        Inches(0.08),
        style["secondary"],
    )

    add_slide_number(
        slide,
        number,
        style["font"],
        style["primary"],
    )


# =========================================================
# ZAMONAVIY
# =========================================================

def design_modern(
    slide,
    title,
    lines,
    number,
    style,
):
    """
    To'q fonli zamonaviy dizayn.
    """

    set_slide_background(
        slide,
        style["bg"],
    )

    # Katta dekorativ doira
    circle = add_shape(
        slide,
        MSO_SHAPE.OVAL,
        Inches(10.8),
        Inches(-1.1),
        Inches(3.2),
        Inches(3.2),
        style["secondary"],
    )

    # Shaffoflikni imkon qadar kamaytirish
    try:
        circle.fill.transparency = 20
    except Exception:
        pass

    # Oltin chiziq
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        Inches(1.42),
        Inches(2.2),
        Inches(0.08),
        style["primary"],
    )

    add_title(
        slide,
        title,
        style["font"],
        style["title_size"],
        style["primary"],
    )

    add_body(
        slide,
        lines,
        style["font"],
        style["body_size"],
        style["text"],
    )

    # Pastki dekor
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(7.25),
        Inches(13.333),
        Inches(0.25),
        style["primary"],
    )

    add_slide_number(
        slide,
        number,
        style["font"],
        style["primary"],
    )


# =========================================================
# MINIMAL
# =========================================================

def design_minimal(
    slide,
    title,
    lines,
    number,
    style,
):
    """
    Juda toza va sodda dizayn.
    """

    set_slide_background(
        slide,
        style["bg"],
    )

    # Kichik nuqta
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        Inches(0.75),
        Inches(0.55),
        Inches(0.18),
        Inches(0.18),
        style["primary"],
    )

    add_title(
        slide,
        title,
        style["font"],
        style["title_size"],
        style["primary"],
    )

    # Ingichka chiziq
    add_shape(
        slide,
        MSO_SHAPE.RECTANGLE,
        Inches(0.8),
        Inches(1.45),
        Inches(11.7),
        Inches(0.025),
        style["secondary"],
    )

    add_body(
        slide,
        lines,
        style["font"],
        style["body_size"],
        style["text"],
    )

    add_slide_number(
        slide,
        number,
        style["font"],
        style["secondary"],
    )


# =========================================================
# KREATIV
# =========================================================

def design_creative(
    slide,
    title,
    lines,
    number,
    style,
):
    """
    Rangli va kreativ talabalar dizayni.
    """

    set_slide_background(
        slide,
        style["bg"],
    )

    # Binafsha yuqori blok
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.55),
        Inches(0.35),
        Inches(12.2),
        Inches(1.15),
        style["primary"],
    )

    add_title(
        slide,
        title,
        style["font"],
        style["title_size"],
        style["white"],
    )

    # O'ng yuqoridagi dekor
    add_shape(
        slide,
        MSO_SHAPE.OVAL,
        Inches(11.75),
        Inches(-0.45),
        Inches(1.6),
        Inches(1.6),
        style["secondary"],
    )

    # Matn uchun oq blok
    body_shape = add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(1.85),
        Inches(11.8),
        Inches(4.7),
        style["white"],
    )

    # Matn blokining chizig'ini olib tashlash
    body_shape.line.fill.background()

    add_body(
        slide,
        lines,
        style["font"],
        style["body_size"],
        style["text"],
    )

    # Pastki dekor
    add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(6.75),
        Inches(1.2),
        Inches(0.18),
        style["secondary"],
    )

    add_slide_number(
        slide,
        number,
        style["font"],
        style["primary"],
    )


# =========================================================
# SLAYD YARATISH
# =========================================================

def make_pptx(
    title,
    body,
    n=10,
    design="🎓 Akademik",
):

    safe_title = (
        title[:60]
        .replace("/", "_")
        .replace("\\", "_")
        .replace(":", "_")
        .replace("*", "_")
        .replace("?", "_")
        .replace('"', "_")
        .replace("<", "_")
        .replace(">", "_")
        .replace("|", "_")
    )

    path = Path("output") / f"{safe_title}.pptx"

    presentation = Presentation()

    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    lines = [
        x.strip()
        for x in body.splitlines()
        if x.strip()
    ]

    # 1–30 oralig'ida
    n = max(
        1,
        min(int(n), 30),
    )

    # Matnni slaydlarga taqsimlash
    per_slide = max(
        1,
        (len(lines) + n - 1) // n,
    )

    selected = DESIGNS.get(
        design,
        DESIGNS["🎓 Akademik"],
    )

    for i in range(n):

        chunk = lines[
            i * per_slide:
            (i + 1) * per_slide
        ]

        if not chunk:
            break

        slide = presentation.slides.add_slide(
            presentation.slide_layouts[6]
        )

        # Tanlangan dizayn
        if design == "🎓 Akademik":

            design_academic(
                slide,
                title if i == 0 else f"{title} — {i + 1}",
                chunk,
                i + 1,
                selected,
            )

        elif design == "💼 Professional":

            design_professional(
                slide,
                title if i == 0 else f"{title} — {i + 1}",
                chunk,
                i + 1,
                selected,
            )

        elif design == "✨ Zamonaviy":

            design_modern(
                slide,
                title if i == 0 else f"{title} — {i + 1}",
                chunk,
                i + 1,
                selected,
            )

        elif design == "🧊 Minimal":

            design_minimal(
                slide,
                title if i == 0 else f"{title} — {i + 1}",
                chunk,
                i + 1,
                selected,
            )

        elif design == "🌈 Kreativ":

            design_creative(
                slide,
                title if i == 0 else f"{title} — {i + 1}",
                chunk,
                i + 1,
                selected,
            )

        else:

            design_academic(
                slide,
                title if i == 0 else f"{title} — {i + 1}",
                chunk,
                i + 1,
                DESIGNS["🎓 Akademik"],
            )

    presentation.save(path)

    return path
    return path
